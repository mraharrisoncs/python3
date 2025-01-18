# needs pip install python-pptx, see python-pptx.readthedocs.io
from pptx import Presentation
import re,os

def extract_text(pptx_path):
    presentation = Presentation(pptx_path)
    extracted_text = ""
    for snum, slide in enumerate(presentation.slides):
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            extracted_text += f"\S{snum + 1}: {notes_text}"
    return extracted_text

def parse_for_timings(text):
    total = 0
    # use regex to find timings in the form "Time: 5 mins" or "Timing 5s"
    matches = re.findall(r'\b[Tt]im(?:e|ing|ings)[\s\.,:;]*?(\d+)\s*(\w*)', text)
    # iterate over matches and convert to minutes
    for num, unit in matches:
        total += int(num) / (60 if unit.startswith("s") else 1)
    return total

def enum_pptx_files(folder):
    for root, dirs, files in os.walk(folder):
        for file in files:
            if file.endswith(".pptx"):
                print(f"Processing {file}")
                yield os.path.join(root, file)

def sum_durations(pptx_files):  
    total = 0
    for file in pptx_files:
        text = extract_text(file)
        duration = parse_for_timings(text)
        print(f"{file}: {duration:.2f} minutes")
        total += duration
    return total

# MAIN PROGRAM
pptx_path = r"C:\Users\alan\OneDrive\#NCCE\PDEs\CP411 Copy\Admin"
total = sum_durations(enum_pptx_files(pptx_path))
print(f"Total duration {total:.2f} minutes")
print(f"or {total//60:.2f} hours {total%60:.2f} minutes")  


